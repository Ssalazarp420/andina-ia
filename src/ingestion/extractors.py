"""
Fase 2 - Extracción por formato.
Cada función recibe la ruta de un archivo y devuelve una lista de segmentos
(etiqueta_de_ubicación, texto), en vez de un único string plano. Esto permite
que la Fase 2 (chunking) y la Fase 5 (citación) sepan de qué página, slide,
hoja o sección exacta vino cada fragmento de texto.
El manejo de errores es obligatorio: un archivo corrupto NO debe tumbar el pipeline completo.
"""

from pathlib import Path

Segment = tuple[str, str]  # (etiqueta_de_ubicacion, texto)


class ExtractionError(Exception):
    """Error controlado al extraer contenido de un documento."""


def extract_pdf(file_path: str) -> list[Segment]:
    """
    Extrae texto de un PDF, página por página. Si el PDF es escaneado (imagen),
    debe hacer fallback a OCR (pendiente: integrar pytesseract / OCR de OCI Vision).
    """
    try:
        import pdfplumber

        segments: list[Segment] = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    segments.append((f"Página {i}", page_text))

        if not segments:
            # TODO: Fase 2 - fallback a OCR cuando el PDF es una imagen escaneada
            raise ExtractionError(
                f"'{file_path}' no devolvió texto extraíble. Posible PDF escaneado: requiere OCR."
            )

        return segments

    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Fallo al extraer PDF '{file_path}': {exc}") from exc


def extract_docx(file_path: str) -> list[Segment]:
    """
    Extrae texto de un Word, preservando títulos/párrafos en orden.
    Usa el último título (Heading) visto como etiqueta de sección para los
    párrafos que le siguen, así se conserva la jerarquía del documento.
    """
    try:
        import docx

        document = docx.Document(file_path)
        segments: list[Segment] = []
        current_section = "Inicio del documento"
        buffer: list[str] = []

        def flush():
            text = "\n".join(buffer).strip()
            if text:
                segments.append((current_section, text))

        for p in document.paragraphs:
            if not p.text.strip():
                continue
            if p.style and p.style.name and p.style.name.lower().startswith("heading"):
                flush()
                buffer.clear()
                current_section = p.text.strip()
            else:
                buffer.append(p.text)

        flush()
        return segments
    except Exception as exc:
        raise ExtractionError(f"Fallo al extraer DOCX '{file_path}': {exc}") from exc


def extract_xlsx(file_path: str) -> list[Segment]:
    """
    Convierte cada hoja en un segmento propio, con las filas como frases
    estructuradas que repiten los encabezados de columna (estrategia definida
    en Fase 2 para datos tabulares).
    """
    try:
        import pandas as pd

        sheets = pd.read_excel(file_path, sheet_name=None)
        segments: list[Segment] = []

        for sheet_name, df in sheets.items():
            rows_as_text = [
                ", ".join(f"{col}: {row[col]}" for col in df.columns)
                for _, row in df.iterrows()
            ]
            sheet_text = "\n".join(rows_as_text).strip()
            if sheet_text:
                segments.append((f"Hoja: {sheet_name}", sheet_text))

        return segments
    except Exception as exc:
        raise ExtractionError(f"Fallo al extraer XLSX '{file_path}': {exc}") from exc


def extract_csv(file_path: str) -> list[Segment]:
    """Misma lógica que XLSX pero para CSV (un único segmento, sin hojas)."""
    try:
        import pandas as pd

        df = pd.read_csv(file_path)
        rows_as_text = [
            ", ".join(f"{col}: {row[col]}" for col in df.columns)
            for _, row in df.iterrows()
        ]
        text = "\n".join(rows_as_text).strip()
        return [("Tabla completa", text)] if text else []
    except Exception as exc:
        raise ExtractionError(f"Fallo al extraer CSV '{file_path}': {exc}") from exc


def extract_pptx(file_path: str) -> list[Segment]:
    """Extrae texto de cada diapositiva (un segmento por slide), incluyendo notas del orador."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        segments: list[Segment] = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text_frame.text)

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text

            combined = "\n".join(slide_text).strip()
            if notes.strip():
                combined = f"{combined}\n[Notas del orador] {notes}".strip()

            if combined:
                segments.append((f"Diapositiva {i}", combined))

        return segments
    except Exception as exc:
        raise ExtractionError(f"Fallo al extraer PPTX '{file_path}': {exc}") from exc


def extract_markdown_or_html(file_path: str) -> list[Segment]:
    """Limpia marcas de Markdown/HTML manteniendo el contenido legible (un único segmento)."""
    try:
        from bs4 import BeautifulSoup

        raw = Path(file_path).read_text(encoding="utf-8", errors="ignore")

        if file_path.lower().endswith((".html", ".htm")):
            text = BeautifulSoup(raw, "html.parser").get_text(separator="\n").strip()
        else:
            text = raw  # Markdown se limpia con más detalle en cleaning.py

        return [("Documento completo", text)] if text.strip() else []
    except Exception as exc:
        raise ExtractionError(f"Fallo al extraer Markdown/HTML '{file_path}': {exc}") from exc


EXTRACTOR_BY_EXTENSION = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".csv": extract_csv,
    ".pptx": extract_pptx,
    ".md": extract_markdown_or_html,
    ".html": extract_markdown_or_html,
    ".htm": extract_markdown_or_html,
}


def extract_document(file_path: str) -> list[Segment]:
    """
    Router principal: elige el extractor correcto según la extensión del archivo.
    Devuelve una lista de segmentos (etiqueta_de_ubicacion, texto).
    """
    extension = Path(file_path).suffix.lower()
    extractor = EXTRACTOR_BY_EXTENSION.get(extension)

    if extractor is None:
        raise ExtractionError(f"Formato no soportado: '{extension}' ({file_path})")

    return extractor(file_path)